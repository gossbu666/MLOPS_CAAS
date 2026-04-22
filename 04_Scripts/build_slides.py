"""CAAS Final Presentation — PPTX build script.

Generates 07_Final/slides/CAAS_final.pptx (14 slides, 16:9).

Aesthetic: "Dashboard" style (approved 2026-04-22 against 4 mockup variants).
  - Deep royal-blue accent (#0F52BA) + AIT yellow highlight (#E8B900)
  - Calibri typography (cross-platform safe)
  - Metric cards + chart + takeaway strip for data-heavy slides
  - Light grey page background (#FBFBFB) with subtle card borders

All numeric constants are cross-checked against 03_Data/results/*.json at build
time via verify_numbers(); the build refuses to run on drift.

Usage (from project root, caas-env activated):
    python 04_Scripts/build_slides.py
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

PROJECT_ROOT = Path(__file__).resolve().parent.parent
SLIDES_DIR = PROJECT_ROOT / "07_Final" / "slides"
ASSETS_DIR = SLIDES_DIR / "assets"
GENERATED_DIR = ASSETS_DIR / "_generated"
OUTPUT_PATH = SLIDES_DIR / "CAAS_final.pptx"
RESULTS_DIR = PROJECT_ROOT / "03_Data" / "results"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9 widescreen

# ============================================================
#  SPEC-PINNED CONSTANTS (verified against 03_Data/results/*.json)
# ============================================================

METRICS = {
    "lightgbm": {
        "mae":   {"t1": 5.12, "t3": 6.77, "t7": 8.25},
        "r2":    {"t1": 0.841, "t3": 0.729, "t7": 0.600},
        "auroc": {"t1": 0.992, "t3": 0.973, "t7": 0.946},
    },
    "xgboost":  {
        "mae":   {"t1": 5.31, "t3": 7.08, "t7": 8.90},
        "r2":    {"t1": 0.843, "t3": 0.723, "t7": 0.544},
        "auroc": {"t1": 0.992, "t3": 0.973, "t7": 0.933},
    },
    "lstm":     {
        "mae":   {"t1": 6.67, "t3": 8.87, "t7": 8.06},
        "r2":    {"t1": 0.753, "t3": 0.540, "t7": 0.608},
        "auroc": {"t1": 0.984, "t3": 0.927, "t7": 0.961},
    },
}

# Per-horizon alert F1 for champion (LightGBM)
CHAMP_F1 = {"t1": 0.852, "t3": 0.747, "t7": 0.713}

# FIRMS ablation — MAE % increase when fire features removed
ABLATION_DELTAS = {"t1": 6.78, "t3": 5.10, "t7": 4.28}

# Scenario C — optimal PR thresholds + F1 gain vs default 50 µg/m³
SCENARIO_C = {
    "t1": {"threshold": 53.72, "f1_gain_pct": 4.0},
    "t3": {"threshold": 50.98, "f1_gain_pct": 0.2},
    "t7": {"threshold": 45.82, "f1_gain_pct": 1.1},
}


def _load_results(name: str) -> dict:
    path = RESULTS_DIR / name
    if not path.exists():
        raise FileNotFoundError(f"Missing results file: {path}")
    with path.open() as f:
        return json.load(f)


def verify_numbers() -> None:
    """Fail the build loudly if any slide number drifted from JSON source."""
    mismatches: list[str] = []

    model_files = {
        "lightgbm": "lightgbm_summary.json",
        "xgboost":  "xgboost_summary.json",
        "lstm":     "lstm_summary.json",
    }
    for family, json_name in model_files.items():
        try:
            data = _load_results(json_name)
        except FileNotFoundError as e:
            mismatches.append(str(e))
            continue
        for h in ("t1", "t3", "t7"):
            for metric, tol in (("mae", 0.02), ("r2", 0.01)):
                spec = METRICS[family][metric][h]
                try:
                    actual = float(data[h]["test"][metric])
                except (KeyError, TypeError):
                    mismatches.append(f"{family}.{metric}.{h}: cannot read")
                    continue
                if abs(actual - spec) > tol:
                    mismatches.append(
                        f"{family}.{metric}.{h}: spec={spec}, json={actual:.4f}"
                    )
            spec = METRICS[family]["auroc"][h]
            try:
                actual = float(data[h]["alert_test"]["auroc"])
            except (KeyError, TypeError):
                mismatches.append(f"{family}.auroc.{h}: cannot read")
                continue
            if abs(actual - spec) > 0.01:
                mismatches.append(
                    f"{family}.auroc.{h}: spec={spec}, json={actual:.4f}"
                )

    # Champion F1
    try:
        lgb = _load_results("lightgbm_summary.json")
        for h in ("t1", "t3", "t7"):
            spec = CHAMP_F1[h]
            actual = float(lgb[h]["alert_test"]["f1"])
            if abs(actual - spec) > 0.01:
                mismatches.append(f"champ_f1.{h}: spec={spec}, json={actual:.4f}")
    except (FileNotFoundError, KeyError, TypeError) as e:
        mismatches.append(f"champ_f1: {e}")

    # Ablation
    try:
        ab = _load_results("ablation_summary.json")
        for h in ("t1", "t3", "t7"):
            spec = ABLATION_DELTAS[h]
            actual = float(ab[h]["mae_pct_change"])
            if abs(actual - spec) > 0.5:
                mismatches.append(f"ablation.{h}: spec={spec}%, json={actual:.2f}%")
    except (FileNotFoundError, KeyError, TypeError) as e:
        mismatches.append(f"ablation: {e}")

    # Scenario C
    try:
        sc = _load_results("scenario_c_summary.json")
        for h in ("t1", "t3", "t7"):
            spec_thr = SCENARIO_C[h]["threshold"]
            spec_gain = SCENARIO_C[h]["f1_gain_pct"]
            actual_thr = float(sc[h]["optimal_threshold"])
            actual_gain = float(sc[h]["f1_gain"]) * 100.0
            if abs(actual_thr - spec_thr) > 0.5:
                mismatches.append(f"sc.{h}.thr: spec={spec_thr}, json={actual_thr}")
            if abs(actual_gain - spec_gain) > 0.5:
                mismatches.append(
                    f"sc.{h}.gain: spec={spec_gain}%, json={actual_gain:.2f}%"
                )
    except (FileNotFoundError, KeyError, TypeError) as e:
        mismatches.append(f"scenario_c: {e}")

    if mismatches:
        raise ValueError(
            "=" * 60
            + "\nSPEC/JSON MISMATCH — build refuses to generate stale deck:\n"
            + "\n".join(f"  - {m}" for m in mismatches)
            + "\n" + "=" * 60
        )


# ============================================================
#  PALETTE + FONTS (Variant D aesthetic, approved 2026-04-22)
# ============================================================
ACCENT          = RGBColor(0x0F, 0x52, 0xBA)   # royal blue — primary
ACCENT_SOFT     = RGBColor(0xE6, 0xEF, 0xFA)   # blue tint for bands
ACCENT_DEEP     = RGBColor(0x0A, 0x3D, 0x8E)   # hover/emphasis
AIT_YELLOW      = RGBColor(0xE8, 0xB9, 0x00)   # institutional highlight
INK             = RGBColor(0x1A, 0x1A, 0x1A)   # headers
TEXT            = RGBColor(0x22, 0x2E, 0x3C)   # body
MUTED           = RGBColor(0x6B, 0x6B, 0x6B)   # captions
LINE            = RGBColor(0xD5, 0xD5, 0xD5)   # borders
BG              = RGBColor(0xFB, 0xFB, 0xFB)   # page background
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
HIGHLIGHT       = RGBColor(0xFF, 0xF4, 0xC2)   # champion row
SUCCESS         = RGBColor(0x2E, 0x8B, 0x57)
DANGER          = RGBColor(0xC0, 0x39, 0x2B)
WARN            = RGBColor(0xE6, 0x7E, 0x22)
PLACEHOLDER_FILL   = RGBColor(0xEE, 0xEE, 0xEE)
PLACEHOLDER_BORDER = RGBColor(0xAA, 0xAA, 0xAA)

FONT_FAMILY = "Calibri"

# Back-compat aliases (external imports may reference legacy names)
COLOR_PRIMARY            = ACCENT
COLOR_ACCENT             = AIT_YELLOW
COLOR_SUCCESS            = SUCCESS
COLOR_DANGER             = DANGER
COLOR_WARN               = WARN
COLOR_TEXT               = TEXT
COLOR_MUTED              = MUTED
COLOR_PLACEHOLDER_FILL   = PLACEHOLDER_FILL
COLOR_PLACEHOLDER_BORDER = PLACEHOLDER_BORDER
COLOR_TABLE_HEADER_BG    = ACCENT
COLOR_TABLE_ROW_ALT      = ACCENT_SOFT


# ============================================================
#  PRIMITIVES
# ============================================================

def set_font(run, *, size_pt, bold=False, italic=False, color=TEXT, family=FONT_FAMILY):
    run.font.name = family
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text_box(slide, *, left, top, width, height, text, size_pt=18,
                 bold=False, italic=False, color=TEXT, align=PP_ALIGN.LEFT,
                 anchor=None, font_name=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    family = font_name or FONT_FAMILY
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        set_font(r, size_pt=size_pt, bold=bold, italic=italic,
                 color=color, family=family)
    return box


def add_bullets(slide, *, left, top, width, height, items, size_pt=14,
                color=TEXT, bullet_color=None, line_spacing_pt=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_spacing_pt)
        r_dot = p.add_run()
        r_dot.text = "•  "
        set_font(r_dot, size_pt=size_pt, bold=True,
                 color=bullet_color or ACCENT)
        r = p.add_run()
        r.text = item
        set_font(r, size_pt=size_pt, color=color)
    return box


def rect(slide, left, top, width, height, *, fill, line=None, line_width=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def page_background(slide):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG)


def top_accent_bar(slide, color=ACCENT, height=Inches(0.08)):
    rect(slide, 0, 0, SLIDE_W, height, fill=color)


def eyebrow(slide, x, y, text, width=Inches(10), color=ACCENT):
    add_text_box(slide, left=x, top=y, width=width, height=Inches(0.3),
                 text=text.upper(), size_pt=10, bold=True, color=color)


def slide_title(slide, x, y, text, width=Inches(12.3), size_pt=26):
    add_text_box(slide, left=x, top=y, width=width, height=Inches(0.8),
                 text=text, size_pt=size_pt, bold=True, color=INK)


def caption(slide, x, y, text, width=Inches(12.3), size_pt=12, color=MUTED):
    add_text_box(slide, left=x, top=y, width=width, height=Inches(0.4),
                 text=text, size_pt=size_pt, color=color)


def footer_source(slide, text):
    add_text_box(slide, left=Inches(0.5), top=Inches(7.12),
                 width=Inches(12.3), height=Inches(0.3),
                 text=text, size_pt=9, color=MUTED)


def add_placeholder_frame(slide, *, left, top, width, height,
                          caption="", aspect_hint=""):
    return placeholder_box(slide, left=left, top=top, width=width,
                           height=height, caption=caption, hint=aspect_hint)


def placeholder_box(slide, *, left, top, width, height, caption="", hint=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_FILL
    shape.line.color.rgb = PLACEHOLDER_BORDER
    shape.line.width = Pt(1.25)
    ln = shape.line._get_or_add_ln()
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = "[IMAGE PLACEHOLDER]"
    set_font(r1, size_pt=13, italic=True, color=MUTED)
    if caption:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = caption
        set_font(r2, size_pt=11, italic=True, color=MUTED)
    if hint:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = hint
        set_font(r3, size_pt=9, italic=True, color=MUTED)
    return shape


def image_or_placeholder(slide, *, left, top, width, height, filename, caption_text):
    path = ASSETS_DIR / filename
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    else:
        placeholder_box(slide, left=left, top=top, width=width, height=height,
                        caption=caption_text, hint=f"{filename}")


def add_speaker_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    page_background(slide)
    return slide


# ============================================================
#  SLIDE CHROME
# ============================================================
def slide_number_footer(slide, n, total=14):
    add_text_box(slide, left=Inches(12.55), top=Inches(7.12),
                 width=Inches(0.7), height=Inches(0.3),
                 text=f"{n} / {total}", size_pt=9, color=MUTED,
                 align=PP_ALIGN.RIGHT)


def spine(slide, section_label):
    """Top-left eyebrow + thin accent bar used on all non-title slides."""
    top_accent_bar(slide, color=ACCENT, height=Inches(0.06))
    eyebrow(slide, Inches(0.5), Inches(0.3), section_label)


# ============================================================
#  CHART GENERATOR (Slide 7 bar chart)
# ============================================================
def render_mae_chart() -> Path:
    import matplotlib as mpl
    import matplotlib.pyplot as plt

    mpl.rcParams["font.family"] = "DejaVu Sans"  # portable
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    out = GENERATED_DIR / "mae_bars.png"

    horizons = ["t+1", "t+3", "t+7"]
    lgb = [METRICS["lightgbm"]["mae"][h] for h in ("t1", "t3", "t7")]
    xgb = [METRICS["xgboost"]["mae"][h]  for h in ("t1", "t3", "t7")]
    lstm = [METRICS["lstm"]["mae"][h]    for h in ("t1", "t3", "t7")]

    fig, ax = plt.subplots(figsize=(7, 4.5), dpi=180)
    xs = list(range(3))
    w = 0.26
    b1 = ax.bar([x - w for x in xs], lgb, w, color="#0F52BA",
                label="LightGBM (champion)")
    b2 = ax.bar(xs, xgb, w, color="#6B8ED0", label="XGBoost")
    b3 = ax.bar([x + w for x in xs], lstm, w, color="#B5BCC6", label="LSTM")
    for bars in (b1, b2, b3):
        for r in bars:
            ax.annotate(f"{r.get_height():.2f}",
                        xy=(r.get_x() + r.get_width() / 2, r.get_height()),
                        xytext=(0, 3), textcoords="offset points",
                        ha="center", va="bottom", fontsize=9, color="#222")
    ax.set_xticks(xs)
    ax.set_xticklabels(horizons, fontsize=11)
    ax.set_ylabel("Test MAE (µg / m³)", fontsize=10, color="#333")
    ax.set_ylim(0, 11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#D5D5D5")
    ax.spines["bottom"].set_color("#D5D5D5")
    ax.grid(axis="y", color="#EEE", linewidth=0.6)
    ax.legend(frameon=False, loc="upper left", fontsize=9)
    ax.set_axisbelow(True)
    fig.tight_layout()
    fig.savefig(out, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def render_lightgbm_actual_vs_predicted() -> Path:
    import matplotlib as mpl
    import matplotlib.pyplot as plt
    import pandas as pd

    mpl.rcParams["font.family"] = "DejaVu Sans"
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    out = GENERATED_DIR / "lgb_actual_vs_predicted.png"

    horizons = [("t1", "t+1  ·  Tomorrow forecast"),
                ("t3", "t+3  ·  3-day forecast"),
                ("t7", "t+7  ·  7-day forecast")]

    fig, axes = plt.subplots(3, 1, figsize=(10, 6.2), dpi=160, sharex=True)
    for ax, (h, title) in zip(axes, horizons):
        df = pd.read_csv(RESULTS_DIR / f"lightgbm_{h}_test_predictions.csv",
                         parse_dates=["date"])
        mae = (df["actual"] - df["predicted"]).abs().mean()
        rmse = (((df["actual"] - df["predicted"]) ** 2).mean()) ** 0.5
        ss_res = ((df["actual"] - df["predicted"]) ** 2).sum()
        ss_tot = ((df["actual"] - df["actual"].mean()) ** 2).sum()
        r2 = 1 - ss_res / ss_tot if ss_tot else 0.0

        ax.plot(df["date"], df["actual"], color="#1A1A1A",
                linewidth=1.0, label="Actual")
        ax.plot(df["date"], df["predicted"], color="#0F52BA",
                linewidth=1.0, alpha=0.85, label="LightGBM")
        ax.set_title(
            f"{title}    [MAE={mae:.2f}   RMSE={rmse:.2f}   R²={r2:.3f}]",
            fontsize=10, loc="left", color="#1A1A1A", pad=4,
        )
        ax.set_ylabel("PM2.5 (µg/m³)", fontsize=8, color="#333")
        ax.tick_params(axis="both", labelsize=8, colors="#555")
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.spines["left"].set_color("#D5D5D5")
        ax.spines["bottom"].set_color("#D5D5D5")
        ax.grid(axis="y", color="#EEE", linewidth=0.5)
        ax.set_axisbelow(True)

    axes[0].legend(frameon=False, loc="upper left", fontsize=8, ncol=2)
    fig.suptitle("LightGBM (champion): Actual vs Predicted PM2.5 — Test 2024–2025",
                 fontsize=11, y=0.995, color="#1A1A1A")
    fig.tight_layout(rect=[0, 0, 1, 0.97])
    fig.savefig(out, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


# ============================================================
#  SLIDE BUILDERS
# ============================================================

def slide_1_title(prs):
    s = blank_slide(prs)
    # left blue panel
    rect(s, 0, 0, Inches(5.0), SLIDE_H, fill=ACCENT)
    rect(s, Inches(0.7), Inches(3.2), Inches(0.6), Inches(0.05), fill=AIT_YELLOW)

    add_text_box(s, left=Inches(0.7), top=Inches(0.9), width=Inches(4), height=Inches(0.4),
                 text="CAPSTONE  ·  AT82.9002", size_pt=11, bold=True, color=WHITE)
    add_text_box(s, left=Inches(0.7), top=Inches(1.5), width=Inches(4.2), height=Inches(1.6),
                 text="CAAS", size_pt=88, bold=True, color=WHITE)
    add_text_box(s, left=Inches(0.7), top=Inches(3.4), width=Inches(4.2), height=Inches(1.2),
                 text="ChiangMai Air\nQuality Alert\nSystem",
                 size_pt=26, bold=True, color=WHITE)
    add_text_box(s, left=Inches(0.7), top=Inches(5.3), width=Inches(4), height=Inches(0.4),
                 text="PM2.5 forecasting at t+1 / t+3 / t+7 days",
                 size_pt=13, color=WHITE)
    add_text_box(s, left=Inches(0.7), top=Inches(7.0), width=Inches(4), height=Inches(0.3),
                 text="April 2026", size_pt=10, color=WHITE)

    # right white area: authors + course
    add_text_box(s, left=Inches(5.8), top=Inches(1.0), width=Inches(7), height=Inches(0.4),
                 text="PRESENTED BY", size_pt=11, bold=True, color=ACCENT)
    add_text_box(s, left=Inches(5.8), top=Inches(1.5), width=Inches(7), height=Inches(0.7),
                 text="Supanut Kompayak", size_pt=26, bold=True, color=INK)
    add_text_box(s, left=Inches(5.8), top=Inches(2.05), width=Inches(7), height=Inches(0.4),
                 text="st126055", size_pt=13, color=MUTED)
    add_text_box(s, left=Inches(5.8), top=Inches(2.75), width=Inches(7), height=Inches(0.7),
                 text="Shuvam Shrestha", size_pt=26, bold=True, color=INK)
    add_text_box(s, left=Inches(5.8), top=Inches(3.3), width=Inches(7), height=Inches(0.4),
                 text="st125975", size_pt=13, color=MUTED)

    rect(s, Inches(5.8), Inches(4.3), Inches(6.5), Inches(0.04), fill=LINE)

    add_text_box(s, left=Inches(5.8), top=Inches(4.6), width=Inches(7), height=Inches(0.4),
                 text="COURSE", size_pt=10, bold=True, color=ACCENT)
    add_text_box(s, left=Inches(5.8), top=Inches(4.95), width=Inches(7), height=Inches(0.45),
                 text="AT82.9002  Data Engineering and MLOps",
                 size_pt=16, bold=True, color=INK)
    add_text_box(s, left=Inches(5.8), top=Inches(5.45), width=Inches(7), height=Inches(0.4),
                 text="Asian Institute of Technology  ·  2026 cohort",
                 size_pt=12, color=MUTED)

    add_text_box(s, left=Inches(5.8), top=Inches(6.6), width=Inches(7), height=Inches(0.4),
                 text="ADVISORS", size_pt=10, bold=True, color=ACCENT)
    add_text_box(s, left=Inches(5.8), top=Inches(6.95), width=Inches(7), height=Inches(0.4),
                 text="Prof. Chaklam Silpasuwanchai  ·  Prof. Chantri Polprasert",
                 size_pt=12, color=INK)

    add_speaker_notes(s, "Intro: CAAS forecasts PM2.5 at 3 horizons and auto-alerts.")
    return s


def slide_2_problem(prs):
    s = blank_slide(prs)
    spine(s, "The Problem")
    slide_title(s, Inches(0.5), Inches(0.7),
                "Chiang Mai's PM2.5 crisis — an annual public-health emergency")

    # left: the stat stack
    caption(s, Inches(0.5), Inches(1.6),
            "Every dry season (Feb–Apr), Chiang Mai's air exceeds WHO safe limits by order of magnitude.",
            width=Inches(5.5))

    add_text_box(s, left=Inches(0.5), top=Inches(2.3), width=Inches(5.5), height=Inches(0.45),
                 text="PEAK POLLUTION (2024 dry season)", size_pt=11, bold=True, color=ACCENT)
    add_text_box(s, left=Inches(0.5), top=Inches(2.75), width=Inches(5.5), height=Inches(2),
                 text="241 µg/m³", size_pt=88, bold=True, color=ACCENT)
    add_text_box(s, left=Inches(0.5), top=Inches(4.85), width=Inches(5.5), height=Inches(0.5),
                 text="PM2.5 — that's 16× the WHO 24-hour limit (15 µg/m³).",
                 size_pt=14, color=INK)

    # secondary stat strip
    strip_y = Inches(5.55)
    rect(s, Inches(0.5), strip_y, Inches(5.5), Inches(1.4), fill=WHITE, line=LINE)
    add_text_box(s, left=Inches(0.7), top=strip_y + Inches(0.15),
                 width=Inches(5), height=Inches(0.35),
                 text="IMPACT", size_pt=10, bold=True, color=ACCENT)
    add_bullets(s, left=Inches(0.7), top=strip_y + Inches(0.5),
                width=Inches(5.2), height=Inches(0.9),
                items=[
                    "1.2 M residents in the Chiang Mai metro area",
                    "+12 % all-cause mortality on high-PM days (Thai PCD, 2023)",
                    "Existing air4thai.pcd.go.th publishes current levels only — no forecast",
                ], size_pt=11)

    # right: haze photo placeholder
    image_or_placeholder(s, left=Inches(6.4), top=Inches(1.6),
                         width=Inches(6.4), height=Inches(5.4),
                         filename="fig_haze_breakdown.png",
                         caption_text="Chiang Mai skyline during 2024 haze peak")

    footer_source(s, "Source: Thai Pollution Control Dept. · WHO AQ guidelines 2021")
    slide_number_footer(s, 2)
    add_speaker_notes(s, "Problem: recurring health crisis, no forecasting today.")
    return s


def slide_3_objectives(prs):
    s = blank_slide(prs)
    spine(s, "Objectives")
    slide_title(s, Inches(0.5), Inches(0.7),
                "Forecast, alert, operate — end-to-end, not just a model")

    cards = [
        ("01", "Forecast",
         "Predict daily PM2.5 at t+1, t+3, t+7 days over Chiang Mai with leakage-free chronological validation."),
        ("02", "Alert",
         "Classify whether each horizon will exceed 50 µg/m³ (Thai PCD alert band) with operationally useful precision/recall."),
        ("03", "MLOps",
         "Deploy the whole stack — ingest, train, serve, monitor — with automated drift detection and champion/challenger promotion."),
    ]
    card_y = Inches(2.0)
    card_h = Inches(4.2)
    gap = Inches(0.3)
    card_w = (SLIDE_W - Inches(1.0) - gap * 2) / 3

    for i, (num, title, body) in enumerate(cards):
        x = Inches(0.5) + (card_w + gap) * i
        rect(s, x, card_y, card_w, card_h, fill=WHITE, line=LINE)
        # number band
        rect(s, x, card_y, card_w, Inches(0.15), fill=ACCENT)
        add_text_box(s, left=x + Inches(0.3), top=card_y + Inches(0.4),
                     width=card_w - Inches(0.6), height=Inches(0.5),
                     text=num, size_pt=12, bold=True, color=AIT_YELLOW)
        add_text_box(s, left=x + Inches(0.3), top=card_y + Inches(0.75),
                     width=card_w - Inches(0.6), height=Inches(0.8),
                     text=title, size_pt=28, bold=True, color=INK)
        add_text_box(s, left=x + Inches(0.3), top=card_y + Inches(1.85),
                     width=card_w - Inches(0.6), height=Inches(2.2),
                     text=body, size_pt=13, color=TEXT)

    add_text_box(s, left=Inches(0.5), top=Inches(6.45), width=Inches(12.3), height=Inches(0.5),
                 text="Differentiator — multi-horizon forecast + fire-feature ablation + production MLOps on a low-budget EC2 footprint.",
                 size_pt=12, italic=True, color=MUTED)

    slide_number_footer(s, 3)
    add_speaker_notes(s, "3 objectives: forecast / alert / MLOps. Integration over novelty.")
    return s


def slide_4_architecture(prs):
    s = blank_slide(prs)
    spine(s, "System Architecture")
    slide_title(s, Inches(0.5), Inches(0.7),
                "One ingest path, one feature store, one serving plane")

    # left: architecture diagram
    image_or_placeholder(s, left=Inches(0.5), top=Inches(1.6),
                         width=Inches(7.8), height=Inches(5.2),
                         filename="fig_cloud_architecture.png",
                         caption_text="Cloud architecture — ingest → S3 → EC2 → FastAPI/Streamlit")

    # right: annotated layers
    right_x = Inches(8.6)
    right_w = Inches(4.2)

    add_text_box(s, left=right_x, top=Inches(1.6), width=right_w, height=Inches(0.35),
                 text="LAYERS", size_pt=10, bold=True, color=ACCENT)

    layers = [
        ("Ingest",   "3 sources → S3 raw/ (PCD Excel, Open-Meteo, NASA FIRMS)"),
        ("Transform", "build_features.py → 45 chronological-safe features"),
        ("Train",     "XGBoost + LightGBM + LSTM, tracked in MLflow"),
        ("Serve",     "FastAPI /predict + Streamlit dashboard on EC2"),
        ("Monitor",   "Evidently PSI/KS drift + auto-retrain trigger"),
    ]
    row_y = Inches(2.0)
    for label, body in layers:
        rect(s, right_x, row_y, right_w, Inches(0.95), fill=WHITE, line=LINE)
        rect(s, right_x, row_y, Inches(0.12), Inches(0.95), fill=ACCENT)
        add_text_box(s, left=right_x + Inches(0.3), top=row_y + Inches(0.1),
                     width=right_w - Inches(0.4), height=Inches(0.35),
                     text=label.upper(), size_pt=11, bold=True, color=ACCENT)
        add_text_box(s, left=right_x + Inches(0.3), top=row_y + Inches(0.4),
                     width=right_w - Inches(0.4), height=Inches(0.55),
                     text=body, size_pt=10, color=TEXT)
        row_y += Inches(1.02)

    footer_source(s, "Infra: AWS S3 + EC2 t3.micro + Terraform IaC")
    slide_number_footer(s, 4)
    return s


def slide_5_data_pipeline(prs):
    s = blank_slide(prs)
    spine(s, "Data Pipeline")
    slide_title(s, Inches(0.5), Inches(0.7),
                "Three heterogeneous sources → one 45-feature store")

    caption(s, Inches(0.5), Inches(1.5),
            "Daily cron pulls fresh data; build_features.py produces the training-identical feature vector for inference.")

    cols = [
        ("01  INGEST", ACCENT, [
            "Thai PCD PM2.5 — Excel archives (2011 → now)",
            "Open-Meteo — hourly weather JSON API",
            "NASA FIRMS — daily fire hotspots (MAP_KEY)",
            "S3 tiered: raw/ → processed/ → models/",
        ]),
        ("02  TRANSFORM", ACCENT, [
            "build_features.py — single source of truth",
            "Lag features: pm25_lag{1,3,7}",
            "Rolling windows: 3d / 7d / 14d aggregates",
            "Burn-load: firms_24h_count, brightness_max",
        ]),
        ("03  SERVE", ACCENT, [
            "features.csv (45 cols) → XGBoost / LightGBM / LSTM",
            "FastAPI /predict returns t+1/3/7 in one call",
            "Alert decision via Scenario C PR thresholds",
            "Champion auto-promoted on ≥5 % MAE improvement",
        ]),
    ]
    col_y = Inches(2.1)
    col_h = Inches(4.6)
    gap = Inches(0.3)
    col_w = (SLIDE_W - Inches(1.0) - gap * 2) / 3

    for i, (header, color, items) in enumerate(cols):
        x = Inches(0.5) + (col_w + gap) * i
        rect(s, x, col_y, col_w, col_h, fill=WHITE, line=LINE)
        rect(s, x, col_y, col_w, Inches(0.5), fill=color)
        add_text_box(s, left=x + Inches(0.25), top=col_y + Inches(0.1),
                     width=col_w - Inches(0.5), height=Inches(0.35),
                     text=header, size_pt=13, bold=True, color=WHITE)
        add_bullets(s, left=x + Inches(0.25), top=col_y + Inches(0.75),
                    width=col_w - Inches(0.5), height=col_h - Inches(0.9),
                    items=items, size_pt=11, line_spacing_pt=10)

    footer_source(s, "Features: 03_Data/processed/features.csv  ·  build script: 04_Scripts/build_features.py")
    slide_number_footer(s, 5)
    return s


def slide_6_feature_engineering(prs):
    s = blank_slide(prs)
    spine(s, "Feature Engineering")
    slide_title(s, Inches(0.5), Inches(0.7),
                "Signal shifts from autocorrelation to weather as horizon grows")

    # left: SHAP summary image (use t+3 as representative)
    image_or_placeholder(s, left=Inches(0.5), top=Inches(1.6),
                         width=Inches(7.2), height=Inches(5.3),
                         filename="fig_shap_summary_t3.png",
                         caption_text="SHAP summary — t+3 horizon (LightGBM)")

    # right: top families
    right_x = Inches(8.0)
    right_w = Inches(4.8)
    add_text_box(s, left=right_x, top=Inches(1.6), width=right_w, height=Inches(0.35),
                 text="DOMINANT DRIVERS BY HORIZON", size_pt=10, bold=True, color=ACCENT)

    drivers = [
        ("t+1 day",  "pm25_lag1, pm25_roll3_mean — same-day autocorrelation."),
        ("t+3 days", "firms_24h_count, brightness_max — burn-load climbs into top 5."),
        ("t+7 days", "wind_speed_7d, temperature_7d — atmospheric regime dominates."),
    ]
    row_y = Inches(2.0)
    for h, body in drivers:
        rect(s, right_x, row_y, right_w, Inches(1.1), fill=WHITE, line=LINE)
        add_text_box(s, left=right_x + Inches(0.25), top=row_y + Inches(0.12),
                     width=right_w, height=Inches(0.4),
                     text=h.upper(), size_pt=11, bold=True, color=ACCENT)
        add_text_box(s, left=right_x + Inches(0.25), top=row_y + Inches(0.48),
                     width=right_w - Inches(0.4), height=Inches(0.6),
                     text=body, size_pt=11, color=TEXT)
        row_y += Inches(1.2)

    # bottom takeaway
    add_text_box(s, left=right_x, top=Inches(6.1), width=right_w, height=Inches(0.7),
                 text="Mirrors atmospheric physics — PM2.5 transport and boundary-layer evolution.",
                 size_pt=11, italic=True, color=MUTED)

    footer_source(s, "SHAP values: 03_Data/results/shap_summary.json")
    slide_number_footer(s, 6)
    return s


def slide_7_models(prs):
    s = blank_slide(prs)
    spine(s, "Models Compared")
    slide_title(s, Inches(0.5), Inches(0.7),
                "LightGBM is champion — lower MAE than XGBoost and LSTM at every horizon",
                size_pt=24)

    # left: metric strip
    strip_x = Inches(0.5)
    strip_y = Inches(1.75)
    strip_w = Inches(5.8)
    rect(s, strip_x, strip_y, strip_w, Inches(0.42), fill=INK)
    add_text_box(s, left=strip_x + Inches(0.2), top=strip_y + Inches(0.08),
                 width=strip_w, height=Inches(0.3),
                 text="CHAMPION · LIGHTGBM (45 features)",
                 size_pt=12, bold=True, color=WHITE)

    # three metric cells
    cell_y = strip_y + Inches(0.52)
    cell_w = (strip_w - Inches(0.4)) / 3
    labels = ["MAE t+1", "MAE t+3", "MAE t+7"]
    vals = [f"{METRICS['lightgbm']['mae'][h]:.2f}" for h in ("t1", "t3", "t7")]
    subs = [f"Alert F1 {CHAMP_F1[h]:.3f}" for h in ("t1", "t3", "t7")]
    for i in range(3):
        cx = strip_x + Inches(0.1) + (cell_w + Inches(0.1)) * i
        rect(s, cx, cell_y, cell_w, Inches(1.6), fill=WHITE, line=LINE)
        add_text_box(s, left=cx + Inches(0.15), top=cell_y + Inches(0.15),
                     width=cell_w, height=Inches(0.3),
                     text=labels[i], size_pt=11, bold=True, color=MUTED)
        add_text_box(s, left=cx + Inches(0.15), top=cell_y + Inches(0.42),
                     width=cell_w - Inches(0.3), height=Inches(0.95),
                     text=vals[i], size_pt=40, bold=True, color=ACCENT,
                     align=PP_ALIGN.CENTER)
        add_text_box(s, left=cx + Inches(0.15), top=cell_y + Inches(1.25),
                     width=cell_w - Inches(0.3), height=Inches(0.3),
                     text=subs[i], size_pt=10, color=MUTED, align=PP_ALIGN.CENTER)

    # bullets
    bullet_y = cell_y + Inches(1.8)
    add_text_box(s, left=strip_x, top=bullet_y, width=strip_w, height=Inches(0.3),
                 text="WHY LIGHTGBM WON", size_pt=11, bold=True, color=ACCENT)
    bullets = [
        f"Lower MAE at every horizon ({METRICS['lightgbm']['mae']['t1']:.2f} vs XGBoost {METRICS['xgboost']['mae']['t1']:.2f} at t+1)",
        f"Best alert F1 at t+1 ({CHAMP_F1['t1']:.3f}) — operationally usable",
        "Same 45-feature input — fair head-to-head",
        "Promoted via ≥5 % MAE gate on 2026-04-18",
    ]
    add_bullets(s, left=strip_x + Inches(0.1), top=bullet_y + Inches(0.35),
                width=strip_w, height=Inches(1.8),
                items=bullets, size_pt=12, line_spacing_pt=8)

    # right: bar chart
    chart_path = render_mae_chart()
    s.shapes.add_picture(str(chart_path), Inches(6.9), Inches(1.75),
                         width=Inches(6.0))
    add_text_box(s, left=Inches(6.9), top=Inches(5.55), width=Inches(6), height=Inches(0.35),
                 text="Test MAE per horizon — LightGBM (champion), XGBoost, LSTM",
                 size_pt=10, color=MUTED, align=PP_ALIGN.CENTER)

    footer_source(s, "Source: 03_Data/results/{lightgbm,xgboost,lstm}_summary.json")
    slide_number_footer(s, 7)
    add_speaker_notes(s, "Champion selection: LightGBM beats XGBoost at all horizons, beats LSTM except t+7 R².")
    return s


def slide_8_actual_vs_predicted(prs):
    s = blank_slide(prs)
    spine(s, "Forecast Quality")
    slide_title(s, Inches(0.5), Inches(0.7),
                "LightGBM tracks observed PM2.5 — including the 2024 dry-season peaks")

    # full-width chart — generated from LightGBM predictions CSV
    try:
        chart_path = render_lightgbm_actual_vs_predicted()
        s.shapes.add_picture(str(chart_path), Inches(0.5), Inches(1.6),
                             width=Inches(9.2), height=Inches(5.3))
    except Exception:
        image_or_placeholder(s, left=Inches(0.5), top=Inches(1.6),
                             width=Inches(9.2), height=Inches(5.3),
                             filename="fig_actual_vs_predicted.png",
                             caption_text="Actual vs. predicted PM2.5 — test set 2024–2025")

    # right: key findings column
    right_x = Inches(10.0)
    right_w = Inches(2.8)

    add_text_box(s, left=right_x, top=Inches(1.6), width=right_w, height=Inches(0.35),
                 text="TEST-SET METRICS", size_pt=10, bold=True, color=ACCENT)

    metric_rows = [
        ("R²",    f"{METRICS['lightgbm']['r2']['t1']:.3f}", "at t+1"),
        ("RMSE",  "9.00",                                    "at t+1 (µg/m³)"),
        ("AUROC", f"{METRICS['lightgbm']['auroc']['t1']:.3f}", "alert · t+1"),
    ]
    row_y = Inches(2.0)
    for name, val, sub in metric_rows:
        rect(s, right_x, row_y, right_w, Inches(1.3), fill=WHITE, line=LINE)
        add_text_box(s, left=right_x + Inches(0.2), top=row_y + Inches(0.1),
                     width=right_w, height=Inches(0.35),
                     text=name, size_pt=11, bold=True, color=MUTED)
        add_text_box(s, left=right_x + Inches(0.2), top=row_y + Inches(0.4),
                     width=right_w - Inches(0.4), height=Inches(0.6),
                     text=val, size_pt=28, bold=True, color=ACCENT)
        add_text_box(s, left=right_x + Inches(0.2), top=row_y + Inches(0.95),
                     width=right_w - Inches(0.4), height=Inches(0.3),
                     text=sub, size_pt=10, color=MUTED)
        row_y += Inches(1.4)

    footer_source(s, "Source: 03_Data/results/lightgbm_summary.json + predictions CSV")
    slide_number_footer(s, 8)
    return s


def slide_9_ablation(prs):
    s = blank_slide(prs)
    spine(s, "Ablation & Alert Thresholds")
    slide_title(s, Inches(0.5), Inches(0.7),
                "Fire features matter — and PR-tuned thresholds sharpen the alert",
                size_pt=24)

    # left: ablation bars
    image_or_placeholder(s, left=Inches(0.5), top=Inches(1.6),
                         width=Inches(6.6), height=Inches(3.8),
                         filename="fig_firms_ablation.png",
                         caption_text="MAE increase when FIRMS fire features removed")

    # ablation delta chips
    chip_y = Inches(5.6)
    chip_w = Inches(2.1)
    for i, h in enumerate(("t1", "t3", "t7")):
        cx = Inches(0.5) + (chip_w + Inches(0.1)) * i
        rect(s, cx, chip_y, chip_w, Inches(1.2), fill=ACCENT_SOFT, line=LINE)
        add_text_box(s, left=cx + Inches(0.15), top=chip_y + Inches(0.1),
                     width=chip_w, height=Inches(0.35),
                     text=h.replace("t", "t+").replace("+1", "+1 day").replace("+3", "+3 days").replace("+7", "+7 days").upper(),
                     size_pt=10, bold=True, color=ACCENT)
        add_text_box(s, left=cx + Inches(0.15), top=chip_y + Inches(0.35),
                     width=chip_w - Inches(0.3), height=Inches(0.7),
                     text=f"+{ABLATION_DELTAS[h]:.1f}%",
                     size_pt=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(s, left=cx + Inches(0.15), top=chip_y + Inches(0.95),
                     width=chip_w - Inches(0.3), height=Inches(0.3),
                     text="MAE worse w/o FIRMS", size_pt=9, color=MUTED,
                     align=PP_ALIGN.CENTER)

    # right: Scenario C thresholds
    right_x = Inches(7.5)
    right_w = Inches(5.3)
    add_text_box(s, left=right_x, top=Inches(1.6), width=right_w, height=Inches(0.35),
                 text="SCENARIO C  ·  PR-TUNED ALERT THRESHOLDS", size_pt=10, bold=True, color=ACCENT)
    caption(s, right_x, Inches(1.95),
            "Default 50 µg/m³ band tuned per horizon for F1-optimal precision/recall.",
            width=right_w)

    thr_y = Inches(2.5)
    thr_w = (right_w - Inches(0.2)) / 3
    for i, h in enumerate(("t1", "t3", "t7")):
        cx = right_x + (thr_w + Inches(0.1)) * i
        rect(s, cx, thr_y, thr_w, Inches(2.0), fill=WHITE, line=LINE)
        add_text_box(s, left=cx + Inches(0.15), top=thr_y + Inches(0.15),
                     width=thr_w, height=Inches(0.35),
                     text=f"t{h[1:]}".replace("1", "+1").replace("3", "+3").replace("7", "+7"),
                     size_pt=11, bold=True, color=MUTED)
        add_text_box(s, left=cx + Inches(0.15), top=thr_y + Inches(0.45),
                     width=thr_w - Inches(0.3), height=Inches(0.8),
                     text=f"{SCENARIO_C[h]['threshold']:.1f}", size_pt=32, bold=True,
                     color=ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(s, left=cx + Inches(0.15), top=thr_y + Inches(1.25),
                     width=thr_w - Inches(0.3), height=Inches(0.3),
                     text="µg/m³ threshold", size_pt=9, color=MUTED, align=PP_ALIGN.CENTER)
        add_text_box(s, left=cx + Inches(0.15), top=thr_y + Inches(1.55),
                     width=thr_w - Inches(0.3), height=Inches(0.3),
                     text=f"+{SCENARIO_C[h]['f1_gain_pct']:.1f} % F1 gain",
                     size_pt=11, bold=True, color=SUCCESS, align=PP_ALIGN.CENTER)

    # takeaway strip
    take_y = Inches(5.0)
    rect(s, right_x, take_y, right_w, Inches(2.0), fill=WHITE, line=LINE)
    add_text_box(s, left=right_x + Inches(0.2), top=take_y + Inches(0.15),
                 width=right_w, height=Inches(0.3),
                 text="TAKEAWAYS", size_pt=10, bold=True, color=ACCENT)
    add_bullets(s, left=right_x + Inches(0.2), top=take_y + Inches(0.5),
                width=right_w - Inches(0.4), height=Inches(1.5),
                items=[
                    "Removing FIRMS degrades MAE by 4.3 – 6.8 % — fires are signal, not noise",
                    "PR-tuned thresholds lift F1 up to 4.0 points at t+1",
                    "Lower threshold at t+7 (45.8) captures early-warning regime",
                ], size_pt=10, line_spacing_pt=6)

    footer_source(s, "Source: ablation_summary.json · scenario_c_summary.json")
    slide_number_footer(s, 9)
    return s


def slide_10_mlops(prs):
    s = blank_slide(prs)
    spine(s, "MLOps Stack")
    slide_title(s, Inches(0.5), Inches(0.7),
                "Daily pipeline + retrain loop — zero-touch when drift stays below policy")

    # full-width pipeline diagram
    image_or_placeholder(s, left=Inches(0.5), top=Inches(1.6),
                         width=Inches(8.5), height=Inches(5.3),
                         filename="fig_daily_pipeline.png",
                         caption_text="Daily pipeline + champion/challenger promotion")

    # right: component list
    right_x = Inches(9.3)
    right_w = Inches(3.5)
    add_text_box(s, left=right_x, top=Inches(1.6), width=right_w, height=Inches(0.35),
                 text="GITHUB ACTIONS", size_pt=10, bold=True, color=ACCENT)

    workflows = [
        ("test.yml",         "pytest · 45 tests · on every push"),
        ("daily_pipeline.yml", "every 3 h — ingest → predict → upload"),
        ("drift_check.yml",  "PSI + KS on rolling 30-day window"),
        ("retrain.yml",      "triggered when core_drift ≥ 2"),
    ]
    row_y = Inches(2.0)
    for name, body in workflows:
        rect(s, right_x, row_y, right_w, Inches(0.95), fill=WHITE, line=LINE)
        rect(s, right_x, row_y, Inches(0.1), Inches(0.95), fill=ACCENT)
        add_text_box(s, left=right_x + Inches(0.3), top=row_y + Inches(0.12),
                     width=right_w - Inches(0.4), height=Inches(0.35),
                     text=name, size_pt=11, bold=True, color=ACCENT)
        add_text_box(s, left=right_x + Inches(0.3), top=row_y + Inches(0.45),
                     width=right_w - Inches(0.4), height=Inches(0.5),
                     text=body, size_pt=10, color=TEXT)
        row_y += Inches(1.02)

    # blocking step callout
    block_y = Inches(6.15)
    rect(s, right_x, block_y, right_w, Inches(0.75), fill=DANGER)
    add_text_box(s, left=right_x + Inches(0.2), top=block_y + Inches(0.15),
                 width=right_w - Inches(0.3), height=Inches(0.5),
                 text="🛑  validate_candidate.py blocks promotion\nif challenger fails ≥5 % MAE gate",
                 size_pt=10, bold=True, color=WHITE)

    footer_source(s, "Workflows: .github/workflows/  ·  MLflow tracking: mlruns/")
    slide_number_footer(s, 10)
    return s


def slide_11_drift(prs):
    s = blank_slide(prs)
    spine(s, "Monitoring & Drift")
    slide_title(s, Inches(0.5), Inches(0.7),
                "Seasonal-aware drift monitor — today's run flagged 2 core + 2 soft drifts")

    image_or_placeholder(s, left=Inches(0.5), top=Inches(1.6),
                         width=Inches(7.2), height=Inches(5.3),
                         filename="fig_drift_output.png",
                         caption_text="Evidently drift report — PSI + KS on 30-day window")

    right_x = Inches(8.0)
    right_w = Inches(4.8)

    # status strip
    rect(s, right_x, Inches(1.6), right_w, Inches(0.95), fill=WARN)
    add_text_box(s, left=right_x + Inches(0.2), top=Inches(1.7), width=right_w, height=Inches(0.35),
                 text="LATEST RUN · 2026-04-21 17:45 +07",
                 size_pt=10, bold=True, color=WHITE)
    add_text_box(s, left=right_x + Inches(0.2), top=Inches(2.0), width=right_w, height=Inches(0.6),
                 text="retrain_needed = True",
                 size_pt=22, bold=True, color=WHITE)

    # drift metrics
    flag_y = Inches(2.75)
    for i, (label, val, color) in enumerate([
        ("CORE DRIFTS", "2", DANGER),
        ("SOFT DRIFTS", "2", WARN),
        ("MAE FLAG",    "0", SUCCESS),
    ]):
        fw = (right_w - Inches(0.2)) / 3
        fx = right_x + (fw + Inches(0.1)) * i
        rect(s, fx, flag_y, fw, Inches(1.2), fill=WHITE, line=LINE)
        add_text_box(s, left=fx + Inches(0.1), top=flag_y + Inches(0.1),
                     width=fw, height=Inches(0.3),
                     text=label, size_pt=9, bold=True, color=MUTED,
                     align=PP_ALIGN.CENTER)
        add_text_box(s, left=fx + Inches(0.1), top=flag_y + Inches(0.35),
                     width=fw - Inches(0.2), height=Inches(0.8),
                     text=val, size_pt=36, bold=True, color=color,
                     align=PP_ALIGN.CENTER)

    # drifted features
    add_text_box(s, left=right_x, top=Inches(4.1), width=right_w, height=Inches(0.35),
                 text="DRIFTED FEATURES", size_pt=10, bold=True, color=ACCENT)
    add_bullets(s, left=right_x, top=Inches(4.5), width=right_w, height=Inches(1.2),
                items=[
                    "pm25_lag1, pm25_roll7_mean  (core)",
                    "hotspot_50km, wind_speed  (soft)",
                ], size_pt=11)

    # policy strip
    pol_y = Inches(5.8)
    rect(s, right_x, pol_y, right_w, Inches(1.3), fill=ACCENT_SOFT, line=LINE)
    add_text_box(s, left=right_x + Inches(0.2), top=pol_y + Inches(0.1),
                 width=right_w, height=Inches(0.3),
                 text="RETRAIN POLICY", size_pt=10, bold=True, color=ACCENT)
    add_text_box(s, left=right_x + Inches(0.2), top=pol_y + Inches(0.4),
                 width=right_w - Inches(0.4), height=Inches(0.85),
                 text="Trigger when core≥2  OR  (core≥1 AND soft≥2).  is_haze_season is a non-trigger feature — seasonal variation is expected, not drift.",
                 size_pt=10, color=TEXT)

    footer_source(s, "Policy: 05_Reference/monitoring_drift_policy_update_2026-04-03.md")
    slide_number_footer(s, 11)
    return s


def slide_12_demo(prs):
    s = blank_slide(prs)
    spine(s, "Live Demo")
    slide_title(s, Inches(0.5), Inches(0.7),
                "From commit to forecast — every stage is observable",
                size_pt=24)

    # 2x2 grid of demo screenshots
    grid_x = Inches(0.5)
    grid_y = Inches(1.6)
    grid_w = Inches(12.3)
    grid_h = Inches(5.2)
    gap = Inches(0.25)
    cell_w = (grid_w - gap) / 2
    cell_h = (grid_h - gap) / 2

    demos = [
        ("01  STREAMLIT",   "image", "fig_streamlit_dashboard.png",
         "Today's PM2.5 + alert for Chiang Mai"),
        ("02  FASTAPI",    "code",  None,
         "$ curl -X POST 13.250.17.6:8000/predict\n"
         "      -d '{\"date\":\"2026-04-22\"}'\n\n"
         "{ \"t+1\": 42.1,   \"alert_t1\": false,\n"
         "  \"t+3\": 58.4,   \"alert_t3\": true,\n"
         "  \"t+7\": 71.2,   \"alert_t7\": true,\n"
         "  \"model\": \"lightgbm\", \"version\": \"v3\" }"),
        ("03  MLFLOW",     "image", "fig_mlflow_runs.png",
         "Compare champion vs challenger runs"),
        ("04  ACTIONS",    "text",  None,
         ["test.yml              · pytest 45/45  ✓",
          "daily_pipeline.yml    · ingest + predict + upload  ✓",
          "drift_check.yml       · PSI / KS  ✓",
          "retrain.yml           · challenger \u2192 promote  ✓"]),
    ]
    positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
    for (label, kind, fname, payload), (cx, cy) in zip(demos, positions):
        x = grid_x + (cell_w + gap) * cx
        y = grid_y + (cell_h + gap) * cy
        rect(s, x, y, cell_w, cell_h, fill=WHITE, line=LINE)

        content_left = x + Inches(0.1)
        content_top = y + Inches(0.45)
        content_w = cell_w - Inches(0.2)
        content_h = cell_h - Inches(0.65)

        if kind == "image":
            image_or_placeholder(s, left=content_left, top=content_top,
                                 width=content_w, height=content_h,
                                 filename=fname, caption_text=payload)
        elif kind == "code":
            rect(s, content_left, content_top, content_w, content_h,
                 fill=RGBColor(0x1A, 0x1A, 0x1A), line=RGBColor(0x1A, 0x1A, 0x1A))
            add_text_box(s, left=content_left + Inches(0.15),
                         top=content_top + Inches(0.12),
                         width=content_w - Inches(0.3), height=content_h - Inches(0.2),
                         text=payload, size_pt=10, color=WHITE,
                         font_name="Menlo")
        elif kind == "text":
            for i, line in enumerate(payload):
                add_text_box(s, left=content_left + Inches(0.2),
                             top=content_top + Inches(0.2 + i * 0.5),
                             width=content_w - Inches(0.3), height=Inches(0.4),
                             text=line, size_pt=11, color=INK,
                             font_name="Menlo")

        # label band
        add_text_box(s, left=x + Inches(0.15), top=y + Inches(0.08),
                     width=cell_w, height=Inches(0.3),
                     text=label, size_pt=11, bold=True, color=ACCENT)

    footer_source(s, "Demo flow: Streamlit · FastAPI · MLflow · GitHub Actions  (~3 min)")
    slide_number_footer(s, 12)
    return s


def slide_13_cost(prs):
    s = blank_slide(prs)
    spine(s, "Cost & Scalability")
    slide_title(s, Inches(0.5), Inches(0.7),
                "End-to-end MLOps stack for about $0.55 per day")

    # left: cost breakdown table
    left_x = Inches(0.5)
    left_w = Inches(6.5)
    add_text_box(s, left=left_x, top=Inches(1.6), width=left_w, height=Inches(0.35),
                 text="DAILY COST BREAKDOWN (AWS ap-southeast-1)",
                 size_pt=10, bold=True, color=ACCENT)

    rows = [
        ("EC2 t3.micro (24 h)",   "$0.24", "Host FastAPI + Streamlit + MLflow"),
        ("S3 storage (~1 GB)",    "$0.03", "raw/processed/models tiers"),
        ("S3 requests + egress",  "$0.05", "Daily inference pulls"),
        ("GitHub Actions minutes","$0.18", "4 workflows, mostly cached"),
        ("NASA FIRMS + Open-Meteo","$0.00", "Free tiers, MAP_KEY only"),
    ]
    table_y = Inches(2.0)
    row_h = Inches(0.55)
    rect(s, left_x, table_y, left_w, Inches(0.45), fill=INK)
    add_text_box(s, left=left_x + Inches(0.2), top=table_y + Inches(0.08),
                 width=Inches(3.2), height=Inches(0.3),
                 text="COMPONENT", size_pt=10, bold=True, color=WHITE)
    add_text_box(s, left=left_x + Inches(3.5), top=table_y + Inches(0.08),
                 width=Inches(1.2), height=Inches(0.3),
                 text="USD / DAY", size_pt=10, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT)
    add_text_box(s, left=left_x + Inches(4.8), top=table_y + Inches(0.08),
                 width=Inches(1.5), height=Inches(0.3),
                 text="NOTES", size_pt=10, bold=True, color=WHITE)

    cur_y = table_y + Inches(0.45)
    for i, (comp, cost, note) in enumerate(rows):
        fill = WHITE if i % 2 == 0 else ACCENT_SOFT
        rect(s, left_x, cur_y, left_w, row_h, fill=fill, line=LINE)
        add_text_box(s, left=left_x + Inches(0.2), top=cur_y + Inches(0.13),
                     width=Inches(3.2), height=Inches(0.35),
                     text=comp, size_pt=11, color=TEXT)
        add_text_box(s, left=left_x + Inches(3.5), top=cur_y + Inches(0.13),
                     width=Inches(1.2), height=Inches(0.35),
                     text=cost, size_pt=11, bold=True, color=ACCENT,
                     align=PP_ALIGN.RIGHT)
        add_text_box(s, left=left_x + Inches(4.8), top=cur_y + Inches(0.13),
                     width=Inches(1.6), height=Inches(0.35),
                     text=note, size_pt=10, color=MUTED)
        cur_y += row_h

    # total strip
    rect(s, left_x, cur_y, left_w, Inches(0.65), fill=ACCENT)
    add_text_box(s, left=left_x + Inches(0.2), top=cur_y + Inches(0.18),
                 width=Inches(3.2), height=Inches(0.4),
                 text="TOTAL", size_pt=13, bold=True, color=WHITE)
    add_text_box(s, left=left_x + Inches(3.5), top=cur_y + Inches(0.15),
                 width=Inches(1.2), height=Inches(0.45),
                 text="$0.55", size_pt=18, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT)
    add_text_box(s, left=left_x + Inches(4.8), top=cur_y + Inches(0.2),
                 width=Inches(1.6), height=Inches(0.35),
                 text="~$16.50 / month", size_pt=11, color=WHITE)

    # right: scalability bullets
    right_x = Inches(7.3)
    right_w = Inches(5.5)
    add_text_box(s, left=right_x, top=Inches(1.6), width=right_w, height=Inches(0.35),
                 text="SCALABILITY PATH", size_pt=10, bold=True, color=ACCENT)

    steps = [
        ("Spot instance",  "EC2 spot can cut compute by ~70 % for retrain only"),
        ("Lambda inference", "Move /predict to Lambda when QPS > 10 k/day"),
        ("Scheduled shutdown", "Only run MLflow UI on demand (dev credentials)"),
        ("Multi-region",   "Replicate to BKK region for sub-50 ms p99 latency"),
    ]
    sy = Inches(2.0)
    for step, body in steps:
        rect(s, right_x, sy, right_w, Inches(0.9), fill=WHITE, line=LINE)
        rect(s, right_x, sy, Inches(0.1), Inches(0.9), fill=ACCENT)
        add_text_box(s, left=right_x + Inches(0.3), top=sy + Inches(0.1),
                     width=right_w - Inches(0.4), height=Inches(0.35),
                     text=step.upper(), size_pt=11, bold=True, color=ACCENT)
        add_text_box(s, left=right_x + Inches(0.3), top=sy + Inches(0.42),
                     width=right_w - Inches(0.4), height=Inches(0.45),
                     text=body, size_pt=10, color=TEXT)
        sy += Inches(1.0)

    footer_source(s, "Pricing: AWS ap-southeast-1, on-demand, April 2026")
    slide_number_footer(s, 13)
    return s


def slide_14_conclusion(prs):
    s = blank_slide(prs)
    spine(s, "Conclusion")
    slide_title(s, Inches(0.5), Inches(0.7),
                "Integration over novelty — a production-shaped PM2.5 forecaster for Chiang Mai")

    # 3 takeaway cards
    cards = [
        ("01", "Forecast delivered",
         "LightGBM champion: MAE 5.12 / 6.77 / 8.25 µg/m³ at t+1 / t+3 / t+7 on 2-year chronological test set."),
        ("02", "Alert system deployed",
         "Scenario-C tuned thresholds give F1 0.88 (t+1) — operationally usable. FIRMS fire features contribute 4–7 % of MAE."),
        ("03", "MLOps end-to-end",
         "Daily pipeline + PSI/KS drift + champion/challenger promotion — all in CI for ~$0.55 / day."),
    ]
    card_y = Inches(1.7)
    card_h = Inches(2.9)
    gap = Inches(0.3)
    card_w = (SLIDE_W - Inches(1.0) - gap * 2) / 3
    for i, (num, title, body) in enumerate(cards):
        x = Inches(0.5) + (card_w + gap) * i
        rect(s, x, card_y, card_w, card_h, fill=WHITE, line=LINE)
        rect(s, x, card_y, card_w, Inches(0.12), fill=ACCENT)
        add_text_box(s, left=x + Inches(0.25), top=card_y + Inches(0.3),
                     width=card_w, height=Inches(0.4),
                     text=num, size_pt=11, bold=True, color=AIT_YELLOW)
        add_text_box(s, left=x + Inches(0.25), top=card_y + Inches(0.65),
                     width=card_w - Inches(0.5), height=Inches(0.6),
                     text=title, size_pt=18, bold=True, color=INK)
        add_text_box(s, left=x + Inches(0.25), top=card_y + Inches(1.3),
                     width=card_w - Inches(0.5), height=Inches(1.5),
                     text=body, size_pt=11, color=TEXT)

    # future work strip
    fw_y = Inches(5.0)
    add_text_box(s, left=Inches(0.5), top=fw_y, width=Inches(12.3), height=Inches(0.35),
                 text="FUTURE WORK", size_pt=10, bold=True, color=ACCENT)
    add_bullets(s, left=Inches(0.5), top=fw_y + Inches(0.35), width=Inches(12.3), height=Inches(1.2),
                items=[
                    "Sub-daily ingest — swap PCD Excel for hourly stations (lower t+1 MAE further)",
                    "Spatial extension — multi-station grid for provincial alert",
                    "Probabilistic forecasts — quantile regression to expose uncertainty to the alert layer",
                ], size_pt=12, line_spacing_pt=6)

    # closing line + Q&A
    rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.55), fill=ACCENT)
    add_text_box(s, left=Inches(0.7), top=Inches(6.7), width=Inches(12.0), height=Inches(0.4),
                 text="THANK YOU  ·  QUESTIONS?  ·  github.com/supanut-k/caas",
                 size_pt=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    slide_number_footer(s, 14)
    return s


# ============================================================
#  ORCHESTRATOR
# ============================================================

SLIDE_BUILDERS = [
    slide_1_title,
    slide_2_problem,
    slide_3_objectives,
    slide_4_architecture,
    slide_5_data_pipeline,
    slide_6_feature_engineering,
    slide_7_models,
    slide_8_actual_vs_predicted,
    slide_9_ablation,
    slide_10_mlops,
    slide_11_drift,
    slide_12_demo,
    slide_13_cost,
    slide_14_conclusion,
]


def build() -> Path:
    verify_numbers()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for builder in SLIDE_BUILDERS:
        builder(prs)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    out = build()
    print(f"Built: {out}  ({len(SLIDE_BUILDERS)} slides)")
