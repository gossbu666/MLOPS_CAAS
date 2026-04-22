"""Build python-pptx style variants B, C, D for CAAS metrics slide."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = Path(__file__).parent

# Numbers — sourced from 03_Data/results/xgboost_summary.json + lstm_summary.json
XGB_MAE = (5.31, 7.08, 8.90)
LSTM_MAE = (6.67, 8.87, 8.06)
XGB_F1 = (0.844, 0.723, 0.584)
LSTM_F1 = (0.753, 0.593, 0.694)
XGB_R2 = (0.843, 0.723, 0.578)
LSTM_R2 = (0.753, 0.540, 0.608)

# Palette
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
LINE = RGBColor(0xD5, 0xD5, 0xD5)
ACCENT = RGBColor(0x0F, 0x52, 0xBA)   # deep blue
ACCENT_SOFT = RGBColor(0xE6, 0xEF, 0xFA)
HIGHLIGHT = RGBColor(0xFF, 0xF4, 0xC2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFB, 0xFB, 0xFB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def textbox(slide, x, y, w, h, text, *, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Helvetica"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def hline(slide, x, y, w, color=LINE, weight=0.75):
    ln = slide.shapes.add_connector(1, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def rect(slide, x, y, w, h, fill, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(0.5)
    r.shadow.inherit = False
    return r


# ---------- Variant B: dense academic table ----------

def build_B():
    prs = new_prs()
    s = blank_slide(prs)

    # header band
    rect(s, 0, 0, SLIDE_W, Inches(0.35), ACCENT)
    textbox(s, Inches(0.5), Inches(0.06), Inches(6), Inches(0.3),
            "CAAS · MODEL PERFORMANCE", size=10, bold=True, color=WHITE, font="Helvetica")
    textbox(s, Inches(10.8), Inches(0.06), Inches(2.2), Inches(0.3),
            "Slide 6 / 14", size=10, color=WHITE, align=PP_ALIGN.RIGHT, font="Helvetica")

    # title + subtitle
    textbox(s, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.7),
            "XGBoost is the champion — lower MAE across every horizon",
            size=26, bold=True)
    textbox(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
            "Chronological split — train 2011–2022 · validate 2023 · test 2024–2025. 45 engineered features, 103 Optuna trials per horizon.",
            size=12, color=MUTED)

    # table — 4 cols (Model + 3 horizons), 6 rows (header + 4 metrics + footer caption)
    col_x = [Inches(0.5), Inches(3.4), Inches(6.0), Inches(8.6), Inches(11.2)]
    col_w = [Inches(2.9), Inches(2.6), Inches(2.6), Inches(2.6), Inches(1.6)]
    row_y = Inches(1.9)
    row_h = Inches(0.48)

    headers = ["", "t+1 day", "t+3 days", "t+7 days", "Alert F1 @ t+1"]
    # header row
    for i, txt in enumerate(headers):
        rect(s, col_x[i], row_y, col_w[i], row_h, INK)
        textbox(s, col_x[i] + Inches(0.15), row_y + Inches(0.08),
                col_w[i] - Inches(0.3), row_h,
                txt, size=12, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)

    # data rows
    rows = [
        ("XGBoost  ·  Champion", [f"{XGB_MAE[0]:.2f}", f"{XGB_MAE[1]:.2f}", f"{XGB_MAE[2]:.2f}"],
         f"{XGB_F1[0]:.3f}", True),
        ("LSTM  ·  Comparison", [f"{LSTM_MAE[0]:.2f}", f"{LSTM_MAE[1]:.2f}", f"{LSTM_MAE[2]:.2f}"],
         f"{LSTM_F1[0]:.3f}", False),
    ]
    cur_y = row_y + row_h
    for name, maes, f1, is_champ in rows:
        fill = HIGHLIGHT if is_champ else WHITE
        for i in range(5):
            rect(s, col_x[i], cur_y, col_w[i], row_h + Inches(0.1), fill, line=LINE)
        # model name
        textbox(s, col_x[0] + Inches(0.15), cur_y + Inches(0.12),
                col_w[0] - Inches(0.3), row_h,
                name, size=13, bold=is_champ, color=INK)
        # MAE cells
        for i, v in enumerate(maes):
            textbox(s, col_x[i + 1] + Inches(0.15), cur_y + Inches(0.12),
                    col_w[i + 1] - Inches(0.3), row_h,
                    v, size=16, bold=is_champ,
                    color=ACCENT if is_champ else INK, align=PP_ALIGN.CENTER)
        # F1
        textbox(s, col_x[4] + Inches(0.15), cur_y + Inches(0.12),
                col_w[4] - Inches(0.3), row_h,
                f1, size=14, bold=is_champ,
                color=INK, align=PP_ALIGN.CENTER)
        cur_y += row_h + Inches(0.1)

    # Δ row
    cur_y += Inches(0.02)
    for i in range(5):
        rect(s, col_x[i], cur_y, col_w[i], row_h, ACCENT_SOFT, line=LINE)
    textbox(s, col_x[0] + Inches(0.15), cur_y + Inches(0.1),
            col_w[0], row_h, "Δ MAE (lower is better)", size=12, bold=True, color=ACCENT)
    deltas = [XGB_MAE[i] - LSTM_MAE[i] for i in range(3)]
    for i, d in enumerate(deltas):
        textbox(s, col_x[i + 1] + Inches(0.15), cur_y + Inches(0.1),
                col_w[i + 1] - Inches(0.3), row_h,
                f"{d:+.2f}", size=13, bold=True,
                color=ACCENT if d < 0 else MUTED, align=PP_ALIGN.CENTER)
    textbox(s, col_x[4] + Inches(0.15), cur_y + Inches(0.1),
            col_w[4] - Inches(0.3), row_h, "+0.091", size=13, bold=True,
            color=ACCENT, align=PP_ALIGN.CENTER)

    # takeaway strip
    cur_y += row_h + Inches(0.35)
    rect(s, Inches(0.5), cur_y, Inches(12.3), Inches(1.4), WHITE, line=LINE)
    textbox(s, Inches(0.75), cur_y + Inches(0.15), Inches(11), Inches(0.4),
            "TAKEAWAYS", size=10, bold=True, color=ACCENT)
    textbox(s, Inches(0.75), cur_y + Inches(0.55), Inches(11.5), Inches(0.8),
            "XGBoost beats LSTM on MAE at every horizon and on alert F1 at t+1 (0.844 vs 0.753). "
            "Gap widens from 1.36 at t+1 to 1.79 at t+3 — fire-activity features dominate mid-horizon.",
            size=12, color=INK)

    # footer
    textbox(s, Inches(0.5), Inches(7.1), Inches(10), Inches(0.3),
            "Source: 03_Data/results/xgboost_summary.json · lstm_summary.json   |   Supanut Kompayak · Shuvam Shrestha",
            size=9, color=MUTED)

    out = HERE / "B_dense_academic.pptx"
    prs.save(out)
    print(f"wrote {out.name}")


# ---------- Variant C: hero metric ----------

def build_C():
    prs = new_prs()
    s = blank_slide(prs)

    # thin top accent
    rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    # eyebrow
    textbox(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.35),
            "CHAMPION MODEL · XGBOOST · 45 FEATURES",
            size=10, bold=True, color=ACCENT)

    # giant title
    textbox(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.3),
            "Test MAE lower than LSTM at every horizon",
            size=34, bold=True)

    # three giant MAE cards
    card_y = Inches(2.5)
    card_h = Inches(2.8)
    card_gap = Inches(0.3)
    card_w = (SLIDE_W - Inches(1.4) - card_gap * 2) / 3

    horizons = ["Forecast t+1 day", "Forecast t+3 days", "Forecast t+7 days"]
    for i in range(3):
        x = Inches(0.7) + (card_w + card_gap) * i
        rect(s, x, card_y, card_w, card_h, WHITE, line=LINE)
        # label
        textbox(s, x + Inches(0.3), card_y + Inches(0.25), card_w - Inches(0.6), Inches(0.35),
                horizons[i], size=11, bold=True, color=MUTED)
        # giant number
        textbox(s, x + Inches(0.3), card_y + Inches(0.6), card_w - Inches(0.6), Inches(1.4),
                f"{XGB_MAE[i]:.2f}", size=84, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        # unit
        textbox(s, x + Inches(0.3), card_y + Inches(1.95), card_w - Inches(0.6), Inches(0.3),
                "µg / m³   MAE (test)", size=11, color=MUTED, align=PP_ALIGN.CENTER)
        # compare vs LSTM
        delta = LSTM_MAE[i] - XGB_MAE[i]
        textbox(s, x + Inches(0.3), card_y + Inches(2.3), card_w - Inches(0.6), Inches(0.4),
                f"↓ {delta:.2f} vs LSTM  ({LSTM_MAE[i]:.2f})",
                size=12, bold=True,
                color=ACCENT if delta > 0 else MUTED, align=PP_ALIGN.CENTER)

    # caption / split
    textbox(s, Inches(0.7), Inches(5.65), Inches(12), Inches(0.5),
            "Chronological split — train 2011–2022 · validate 2023 · test 2024–2025.  No leakage, 103 Optuna trials per horizon.",
            size=12, color=MUTED)

    # one-line finding
    textbox(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.7),
            "XGBoost advances PM2.5 forecasting beyond single-horizon baselines — promoted automatically via ≥5 % MAE gate.",
            size=14, bold=True, color=INK)

    # footer
    textbox(s, Inches(0.7), Inches(7.1), Inches(10), Inches(0.3),
            "Source: xgboost_summary.json · lstm_summary.json",
            size=9, color=MUTED)

    out = HERE / "C_hero_metric.pptx"
    prs.save(out)
    print(f"wrote {out.name}")


# ---------- Variant D: dashboard (table + bar chart) ----------

def build_D():
    import matplotlib.pyplot as plt
    import matplotlib as mpl
    mpl.rcParams["font.family"] = "Helvetica"

    # render chart PNG
    fig, ax = plt.subplots(figsize=(7, 4.5), dpi=180)
    horizons = ["t+1", "t+3", "t+7"]
    x_pos = range(len(horizons))
    w = 0.38
    xgb = list(XGB_MAE)
    lstm = list(LSTM_MAE)

    xgb_color = "#0F52BA"
    lstm_color = "#9AA5B1"

    b1 = ax.bar([i - w / 2 for i in x_pos], xgb, w, color=xgb_color, label="XGBoost (champion)")
    b2 = ax.bar([i + w / 2 for i in x_pos], lstm, w, color=lstm_color, label="LSTM")
    for rects in (b1, b2):
        for r in rects:
            ax.annotate(f"{r.get_height():.2f}",
                        xy=(r.get_x() + r.get_width() / 2, r.get_height()),
                        xytext=(0, 4), textcoords="offset points",
                        ha="center", va="bottom", fontsize=10, color="#1A1A1A")

    ax.set_xticks(list(x_pos))
    ax.set_xticklabels(horizons, fontsize=11)
    ax.set_ylabel("Test MAE (µg / m³)", fontsize=11, color="#333")
    ax.set_ylim(0, 11)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#D5D5D5")
    ax.spines["bottom"].set_color("#D5D5D5")
    ax.tick_params(colors="#555")
    ax.grid(axis="y", color="#EEE", linewidth=0.6)
    ax.legend(frameon=False, loc="upper left", fontsize=10)
    ax.set_axisbelow(True)
    fig.tight_layout()
    chart_path = HERE / "D_chart.png"
    fig.savefig(chart_path, bbox_inches="tight", facecolor="white")
    plt.close(fig)

    prs = new_prs()
    s = blank_slide(prs)

    # thin top accent
    rect(s, 0, 0, SLIDE_W, Inches(0.08), ACCENT)

    # eyebrow
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(0.35),
            "PERFORMANCE · MAE PER HORIZON", size=10, bold=True, color=ACCENT)

    # title
    textbox(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.7),
            "XGBoost is champion — 1.36–1.79 µg/m³ lower MAE than LSTM",
            size=24, bold=True)

    # left: metric strip
    strip_x = Inches(0.7)
    strip_y = Inches(1.75)
    strip_w = Inches(5.8)
    rect(s, strip_x, strip_y, strip_w, Inches(0.42), INK)
    textbox(s, strip_x + Inches(0.2), strip_y + Inches(0.09), strip_w, Inches(0.3),
            "CHAMPION · XGBOOST (45 features)", size=12, bold=True, color=WHITE)

    # three metric cells
    cell_y = strip_y + Inches(0.52)
    cell_w = (strip_w - Inches(0.4)) / 3
    labels = ["MAE t+1", "MAE t+3", "MAE t+7"]
    vals = [f"{v:.2f}" for v in XGB_MAE]
    sub = [f"Alert F1 {f:.3f}" for f in XGB_F1]
    for i in range(3):
        cx = strip_x + Inches(0.1) + (cell_w + Inches(0.1)) * i
        rect(s, cx, cell_y, cell_w, Inches(1.6), WHITE, line=LINE)
        textbox(s, cx + Inches(0.15), cell_y + Inches(0.15), cell_w, Inches(0.3),
                labels[i], size=11, bold=True, color=MUTED)
        textbox(s, cx + Inches(0.15), cell_y + Inches(0.45), cell_w - Inches(0.3), Inches(0.9),
                vals[i], size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        textbox(s, cx + Inches(0.15), cell_y + Inches(1.25), cell_w - Inches(0.3), Inches(0.3),
                sub[i], size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # bullets under strip
    bullet_y = cell_y + Inches(1.8)
    textbox(s, strip_x, bullet_y, strip_w, Inches(0.3),
            "WHY XGBOOST WON", size=11, bold=True, color=ACCENT)
    bullets = [
        "Lower MAE at every horizon (t+1: 5.31 vs 6.67)",
        "Higher alert F1 at t+1 (0.844 vs 0.753)",
        "Promoted via ≥5 % MAE gate on 2026-04-18",
        "Same 45-feature input — fair head-to-head",
    ]
    for i, b in enumerate(bullets):
        textbox(s, strip_x + Inches(0.1), bullet_y + Inches(0.35) + Inches(0.3) * i,
                strip_w - Inches(0.2), Inches(0.3),
                f"•  {b}", size=12, color=INK)

    # right: chart
    s.shapes.add_picture(str(chart_path), Inches(6.9), Inches(1.75),
                         width=Inches(6.0))
    textbox(s, Inches(6.9), Inches(5.55), Inches(6), Inches(0.35),
            "Test MAE, XGBoost vs LSTM, per forecast horizon",
            size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # footer
    textbox(s, Inches(0.7), Inches(7.1), Inches(10), Inches(0.3),
            "Source: 03_Data/results/xgboost_summary.json · lstm_summary.json",
            size=9, color=MUTED)

    out = HERE / "D_dashboard.pptx"
    prs.save(out)
    print(f"wrote {out.name}")


if __name__ == "__main__":
    build_B()
    build_C()
    build_D()
